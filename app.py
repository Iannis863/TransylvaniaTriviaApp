import sys
import os

# --- 1. PYTHON 3.13 & FFmpeg PATH STRIKE ---
current_dir = os.path.dirname(os.path.abspath(__file__))
ffmpeg_dir = os.path.join(current_dir, "ffmpeg")
os.environ["PATH"] += os.pathsep + ffmpeg_dir

try:
    import audioop
except ImportError:
    try:
        import audioop_lts as audioop

        sys.modules["audioop"] = audioop
    except ImportError:
        pass

import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
from pydub import AudioSegment
import glob
import io
import re

# Anchor pydub to the folder
ffmpeg_exe = os.path.join(ffmpeg_dir, "ffmpeg.exe")
if os.path.exists(ffmpeg_exe):
    AudioSegment.converter = ffmpeg_exe
    AudioSegment.ffprobe = os.path.join(ffmpeg_dir, "ffprobe.exe")
else:
    st.error(f"❌ FFmpeg NOT FOUND at: {ffmpeg_exe}")

# --- PAGE CONFIG ---
st.set_page_config(page_title="Transylvania Trivia Command", layout="wide", page_icon="🧛")

# --- FOLDER SETUP ---
GRADING_FOLDER = "Grading Sheets"
MASTER_FILE_PATH = os.path.join(GRADING_FOLDER, "Grading Sheet Overall.xlsx")
if not os.path.exists(GRADING_FOLDER):
    os.makedirs(GRADING_FOLDER)


# --- HELPERS ---
def get_ordinal(n):
    if 11 <= (n % 100) <= 13: return f"{n}th"
    return f"{n}" + {1: 'st', 2: 'nd', 3: 'rd'}.get(n % 10, 'th')


def get_dynamic_font_size(text, is_qa=False):
    length = len(str(text))
    if is_qa: return 64 if length < 15 else (52 if length < 30 else 42)
    return 58 if length < 30 else (48 if length < 70 else 38)


# --- 1. THE ADVANCED SCORING ENGINE ---
def run_calculation(df):
    if df.empty: return df
    numeric_cols = ['R1', 'R2', 'R3', 'R4', 'R5', 'Pariu (±)']
    for col in numeric_cols:
        df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0.00)

    def row_math(row):
        rounds = ["R1", "R2", "R3", "R4", "R5"]
        score_sum = 0.00
        j_choice = str(row['Joker'])
        for r in rounds:
            val = float(row[r])
            score_sum += (val * 2) if j_choice == r else val
        return score_sum + float(row['Pariu (±)'])

    df['Total'] = df.apply(row_math, axis=1)
    return df


def manage_scores():
    st.header("🏆 Trivia Scoring Matrix")
    COLS = ['Echipă', 'R1', 'R2', 'R3', 'R4', 'R5', 'Joker', 'Pariu (±)', 'Total']
    if 'teams' not in st.session_state:
        st.session_state.teams = pd.DataFrame(columns=COLS)

    with st.expander("➕ Adaugă Echipă Nouă", expanded=True):
        t_name = st.text_input("Nume Echipă", key="new_team_name_input")
        if st.button("Înregistrează Echipa"):
            if t_name:
                new_t = pd.DataFrame(
                    [{'Echipă': t_name, 'R1': 0.00, 'R2': 0.00, 'R3': 0.00, 'R4': 0.00, 'R5': 0.00, 'Joker': "Niciuna",
                      'Pariu (±)': 0.00, 'Total': 0.00}])
                st.session_state.teams = pd.concat([st.session_state.teams, new_t], ignore_index=True)
                st.rerun()

    st.markdown("---")
    col_input, col_reset, col_push = st.columns([1, 1, 1])
    with col_input:
        edition = st.number_input("Ediția Trivia", min_value=1, step=1, value=1, label_visibility="collapsed")
    with col_reset:
        if st.button("🧹 Reset / Șterge Tot", use_container_width=True):
            st.session_state.teams = pd.DataFrame(columns=COLS)
            if 'reveal_step' in st.session_state: st.session_state.reveal_step = 0
            st.rerun()
    with col_push:
        if not st.session_state.teams.empty:
            if st.button(f"🚀 Push to Ed. {edition}", use_container_width=True):
                # Individual Logic and Master Sheet Logic Integrated Here
                edition_str = get_ordinal(edition)
                session_path = os.path.join(GRADING_FOLDER, f"Grading Sheet {edition_str} edition.xlsx")
                output_session = io.BytesIO()
                with pd.ExcelWriter(output_session, engine='xlsxwriter') as writer:
                    export_df = st.session_state.teams.sort_values(by='Total', ascending=False)
                    export_df.to_excel(writer, index=False, sheet_name='Final Scores')
                with open(session_path, "wb") as f:
                    f.write(output_session.getvalue())

                # 2. UPDATE MASTER OVERALL SHEET
                edition_col = f"Ediția {edition}"
                if os.path.exists(MASTER_FILE_PATH):
                    df_master = pd.read_excel(MASTER_FILE_PATH, engine='openpyxl')
                else:
                    df_master = pd.DataFrame(columns=['Numele echipei', 'Scor Final'])

                current_scores = st.session_state.teams[['Echipă', 'Total']].copy()
                current_scores = current_scores.rename(columns={'Echipă': 'Numele echipei', 'Total': edition_col})

                # Avoid duplicate columns if pushing the same edition twice
                if edition_col in df_master.columns:
                    df_master = df_master.drop(columns=[edition_col])

                df_master = pd.merge(df_master, current_scores, on='Numele echipei', how='outer')

                # Re-sort columns: Name, then Editions in order, then Final Score
                cols = df_master.columns.tolist()
                edition_cols = [c for c in cols if "Ediția" in c]
                edition_cols.sort(key=lambda x: int(re.search(r'\d+', x).group()))
                df_master = df_master[['Numele echipei'] + edition_cols + ['Scor Final']]

                # Fill gaps and calculate total across all editions
                for col in edition_cols:
                    df_master[col] = df_master[col].fillna("-")

                def calculate_total(row):
                    total_val = 0
                    for col in edition_cols:
                        val = row[col]
                        if isinstance(val, (int, float)):
                            total_val += val
                    return total_val

                df_master['Scor Final'] = df_master.apply(calculate_total, axis=1)
                df_master = df_master.sort_values(by='Scor Final', ascending=False)

                try:
                    with pd.ExcelWriter(MASTER_FILE_PATH, engine='xlsxwriter') as writer:
                        df_master.to_excel(writer, index=False, sheet_name='Overall Ranking')
                    st.toast(f"✅ Master Sheet Updated", icon='🧛')
                except PermissionError:
                    st.error(f"❌ ACCES REFUZAT: Închide Excel-ul și apasă PUSH din nou!")

    num_cfg = st.column_config.NumberColumn(format="%.2f", step=0.01)
    edited_df = st.data_editor(
        st.session_state.teams,
        column_config={
            "Echipă": st.column_config.TextColumn("Echipă", width="medium"),
            "Joker": st.column_config.SelectboxColumn("Joker", options=["Niciuna", "R1", "R2", "R3", "R4", "R5"]),
            "Pariu (±)": num_cfg,
            "Total": st.column_config.NumberColumn("Total", disabled=True, format="%.2f"),
            "R1": num_cfg, "R2": num_cfg, "R3": num_cfg, "R4": num_cfg, "R5": num_cfg,
        },
        hide_index=False, use_container_width=True, num_rows="dynamic"
    )
    if not edited_df.equals(st.session_state.teams):
        st.session_state.teams = run_calculation(edited_df)
        st.rerun()


# --- 2. FULLSCREEN REVEAL ENGINE ---
def show_leaderboard_reveal():
    st.header("📽️ Reveal Clasament")
    if st.session_state.teams.empty:
        st.warning("Adaugă echipe în Scoring Matrix pentru a începe.")
        return

    leaderboard = st.session_state.teams[['Echipă', 'Total']].sort_values(by='Total', ascending=True).reset_index(
        drop=True)
    total_teams = len(leaderboard)

    if 'reveal_step' not in st.session_state: st.session_state.reveal_step = 0

    col1, col2 = st.columns([1, 5])
    if col1.button("⏪ Reset"):
        st.session_state.reveal_step = 0
        st.rerun()

    step = st.session_state.reveal_step
    if step >= total_teams:
        st.success("Reveal Complet!")
        return

    current_team = leaderboard.iloc[step]
    rank = total_teams - step
    name_color = "#ffffff"
    if rank == 1:
        name_color = "#FFD700"
    elif rank == 2:
        name_color = "#C0C0C0"
    elif rank == 3:
        name_color = "#CD7F32"

    st.markdown(f"""
        <style>
            .reveal-container {{
                background-color: #000000; padding: 50px 80px; border-radius: 20px; 
                border: 8px solid #00004a; text-align: center; min-height: 55vh; 
                display: flex; flex-direction: column; justify-content: center; align-items: center; 
                box-shadow: 0 0 40px #000080;
            }}
        </style>
        <div class="reveal-container">
            <h2 style="color: #666666; font-size: 30px; margin: 0; font-family: 'Chromium One', sans-serif;">LOCUL</h2>
            <h1 style="color: #ffffff; font-size: 110px; margin: 0; line-height: 1; font-family: 'Chromium One', sans-serif;">{rank}</h1>
            <div style="width: 250px; height: 3px; background-color: #000080; margin: 20px 0;"></div>
            <h2 style="color: {name_color}; font-size: 80px; font-weight: bold; margin: 5px 0;">{current_team['Echipă']}</h2>
            <h3 style="color: #ffffff; font-size: 50px; margin: 0;">{current_team['Total']:.2f}</h3>
        </div>
    """, unsafe_allow_html=True)

    if st.button("Următorul Loc", use_container_width=True):
        st.session_state.reveal_step += 1
        st.rerun()


# --- 3. STYLE & IMAGE HELPERS ---
def place_smart_scaled_image(slide, img_path, target_center_x, target_center_y, max_w=Inches(5.5), max_h=Inches(4.5)):
    exts = ['.jpg', '.jpeg', '.png', '.avif', '.webp']
    final_path = img_path if img_path and os.path.exists(img_path) else None
    if not final_path and img_path:
        base = os.path.splitext(img_path)[0]
        for ext in exts:
            if os.path.exists(base + ext): final_path = base + ext; break
    if not final_path: return

    if os.path.splitext(final_path)[1].lower() in ['.webp', '.avif', '.mpo']:
        img = Image.open(final_path)
        img_io = io.BytesIO()
        img.save(img_io, format='PNG')
        img_io.seek(0)
        img_to_add = img_io
    else:
        img_to_add = final_path

    try:
        pic = slide.shapes.add_picture(img_to_add, 0, 0)
    except ValueError as e:
        raise ValueError(f"Unsupported image format for file '{final_path}': {e}") from e
    ratio = min(max_w / pic.width, max_h / pic.height)
    new_w, new_h = int(pic.width * ratio), int(pic.height * ratio)
    pic.width, pic.height, pic.left, pic.top = new_w, new_h, int(target_center_x - (new_w / 2)), int(
        target_center_y - (new_h / 2))


def add_styled_text(slide, text, left, top, width, height, font_size=None, is_qa=False, bold=False,
                    alignment=PP_ALIGN.CENTER, font_name=None, color_rgb=(255, 255, 255), force_single_line=True):
    if not text or str(text).strip() == "": return
    if font_size is None: font_size = get_dynamic_font_size(text, is_qa)
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap, tf.vertical_anchor = not force_single_line, MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text, p.alignment = str(text), alignment
    if len(p.runs) > 0:
        run = p.runs[0]
        if force_single_line:
            ratio = 1.4 if font_name == "Chromium One" else 1.9
            cur_size = font_size
            while (len(str(text)) * (cur_size / ratio)) > (width / 914400 * 72):
                cur_size -= 5
                if cur_size < 18: break
            font_size = cur_size
        run.font.size, run.font.bold = Pt(font_size), bold
        if font_name: run.font.name = font_name
        run.font.color.rgb = RGBColor(*color_rgb)


# --- 4. CORE GENERATOR ---
def generate_trivia_slides(df):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    def add_bg(slide, path=None):
        bg = path if path else os.path.join("images", "background.jpeg")
        if os.path.exists(bg): slide.shapes.add_picture(bg, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def add_styled_slide(t="", s=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        if t: add_styled_text(slide, t, Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(3), font_size=150,
                              font_name="Chromium One", color_rgb=(255, 215, 0))
        if s: add_styled_text(slide, s, Inches(0.5), Inches(4), prs.slide_width - Inches(1), Inches(3), font_size=125,
                              font_name="Gladiola", color_rgb=(255, 255, 255))
        return slide

    for img in ["First Slide.jpeg", "Second Slide.jpeg", "Third Slide.jpeg", "Fourth Slide.jpeg"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, os.path.join("images", "Slides", img))

    for r in sorted([r for r in df['Round'].unique() if r <= 5]):
        r_data = df[df['Round'] == r]
        base_name = r_data['Round_Name'].iloc[0]
        add_styled_slide(f"Runda {r}", base_name)

        # --- QUESTION SLIDES LOOP ---
        for i, (_, row) in enumerate(r_data.iterrows(), 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_styled_text(slide, row['Round_Name'] if r == 2 else base_name, Inches(0.5), Inches(0.2), Inches(5),
                            Inches(0.8), font_size=44, alignment=PP_ALIGN.LEFT, font_name="Gladiola",
                            color_rgb=(255, 215, 0))
            add_styled_text(slide, f"Întrebarea {i}", Inches(7.8), Inches(0.2), Inches(5), Inches(0.8), font_size=44,
                            alignment=PP_ALIGN.RIGHT, font_name="Gladiola", color_rgb=(255, 215, 0))

            img_n = str(row.get('Picture')) if pd.notna(row.get('Picture')) else None
            path = os.path.join("images", f"Round {r}", img_n) if img_n else None

            # AUDIO LOGIC (Always runs for Round 3)
            if r == 3:
                a_base = str(row['Question'])
                q_audio_folder = os.path.join("audio", "questions")
                a_exts = ['', '.mp3', '.wav', '.mpeg', '.m4a', '.mp4']
                found_q = next((os.path.join(q_audio_folder, os.path.splitext(a_base)[0] + e) for e in a_exts if
                                os.path.exists(os.path.join(q_audio_folder, os.path.splitext(a_base)[0] + e))), None)
                if found_q:
                    final_q_audio = found_q
                    if not found_q.lower().endswith('.mp3'):
                        try:
                            conv_q_path = os.path.splitext(found_q)[0] + "_converted.mp3"
                            if not os.path.exists(conv_q_path): AudioSegment.from_file(found_q).export(conv_q_path,
                                                                                                       format="mp3")
                            final_q_audio = conv_q_path
                        except Exception:
                            pass
                    slide.shapes.add_movie(os.path.abspath(final_q_audio), Inches(-1), Inches(-1), Inches(0.5),
                                           Inches(0.5)).name = "TriviaAudio"

            # QUESTION SLIDE LAYOUT LOGIC
            if r == 2:
                # 2 Images Side-by-Side
                img_q_val = str(row.get('Question'))
                img_path_q = os.path.join("images", f"Round {r}",
                                          img_q_val) if img_q_val and img_q_val.lower() != 'nan' else None

                if img_path_q:
                    place_smart_scaled_image(slide, img_path_q, prs.slide_width * 0.25,
                                             prs.slide_height / 2 + Inches(0.4), max_w=Inches(6.0), max_h=Inches(5.0))
                if path:
                    place_smart_scaled_image(slide, path, prs.slide_width * 0.75, prs.slide_height / 2 + Inches(0.4),
                                             max_w=Inches(6.0), max_h=Inches(5.0))

            elif r == 3:
                if "Continuă Versul" in base_name:
                    # Lyric on Left, Image on Right
                    ans = str(row.get('Answer', ''))
                    parts = ans.split("-", 1) if "-" in ans else [ans, ""]
                    lyric_q = parts[0].strip()

                    add_styled_text(slide, lyric_q, Inches(0.5), Inches(1.5), prs.slide_width / 2 - Inches(0.5), Inches(5),
                                    bold=True, is_qa=True, font_name="Gladiola", force_single_line=False)
                    if path:
                        place_smart_scaled_image(slide, path, prs.slide_width * 0.75, prs.slide_height / 2 + Inches(0.4),
                                                 max_w=Inches(5.5), max_h=Inches(5.0))
                else:
                    if path:
                        place_smart_scaled_image(slide, path, prs.slide_width / 2, prs.slide_height / 2 + Inches(0.4),
                                                 max_w=Inches(9.0), max_h=Inches(5.0))

            else:
                # Standard Logic for R1, R4, R5
                add_styled_text(slide, row['Question'], Inches(0.5), Inches(1.5), Inches(6), Inches(5), bold=True,
                                is_qa=True, font_name="Gladiola", force_single_line=False)
                if path: place_smart_scaled_image(slide, path, Inches(10), prs.slide_height / 2 + Inches(0.5),
                                                  max_w=Inches(5.0), max_h=Inches(4.5))

        add_styled_slide(f"Runda {r}", "Răspunsuri")

        # --- ANSWER SLIDES LOOP ---
        for i, (_, row) in enumerate(r_data.iterrows(), 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_styled_text(slide, row['Round_Name'] if r == 2 else base_name, Inches(0.5), Inches(0.2), Inches(5),
                            Inches(0.8), font_size=44, alignment=PP_ALIGN.LEFT, font_name="Gladiola",
                            color_rgb=(255, 215, 0))
            add_styled_text(slide, f"Răspuns {i}", Inches(7.8), Inches(0.2), Inches(5), Inches(0.8), font_size=44,
                            alignment=PP_ALIGN.RIGHT, font_name="Gladiola", color_rgb=(255, 215, 0))

            img_n = str(row.get('Picture')) if pd.notna(row.get('Picture')) else None
            path = os.path.join("images", f"Round {r}", img_n) if img_n else None

            # ANSWER SLIDE LAYOUT LOGIC
            # --- ADAPTIVE ROUND 3 ANSWER LAYOUT ---
            if r == 3:
                # Lyric Continuation on Left, Image on Right
                if "Continuă Versul" in base_name:
                    ans = str(row.get('Answer', ''))
                    parts = ans.split("-", 1) if "-" in ans else ["", ans]
                    lyric_a = parts[1].strip() if len(parts) > 1 else ans.strip()
                    add_styled_text(slide, lyric_a, Inches(0.5), Inches(1.5), prs.slide_width / 2 - Inches(0.5),
                                    Inches(5), bold=True, is_qa=True, font_name="Gladiola", force_single_line=False)
                    if path:
                        place_smart_scaled_image(slide, path, prs.slide_width * 0.75,
                                                 prs.slide_height / 2 + Inches(0.4), max_w=Inches(5.5),
                                                 max_h=Inches(5.0))

                # Classic: 3-Zone Layout (Song - Image - Artist)
                else:
                    cw, vt = prs.slide_width / 3, (prs.slide_height - Inches(2.5)) / 2
                    ans = str(row['Answer'])
                    parts = ans.split("-", 1) if "-" in ans else [ans, ""]
                    # Left Zone: Song Name
                    add_styled_text(slide, parts[0].strip(), Inches(0.2), vt, cw - Inches(0.4), Inches(2.5),
                                    font_size=48, bold=True, font_name="Gladiola", force_single_line=False)
                    # Middle Zone: Image
                    if path:
                        place_smart_scaled_image(slide, path, prs.slide_width / 2, prs.slide_height / 2,
                                                 max_w=cw + Inches(1.5), max_h=Inches(5.8))
                    # Right Zone: Artist Name
                    if len(parts) > 1:
                        add_styled_text(slide, parts[1].strip(), prs.slide_width - cw + Inches(0.2), vt,
                                        cw - Inches(0.4), Inches(2.5), font_size=48, bold=True, font_name="Gladiola",
                                        force_single_line=False)

                # --- AUDIO ENGINE (SHARED BY BOTH MODES) ---
                a_base = str(row['Question'])
                ans_audio_folder = os.path.join("audio", "answers")
                a_exts = ['', '.mp3', '.wav', '.mpeg', '.m4a', '.mp4']
                found_ans = next((os.path.join(ans_audio_folder, os.path.splitext(a_base)[0] + e) for e in a_exts if
                                  os.path.exists(os.path.join(ans_audio_folder, os.path.splitext(a_base)[0] + e))),
                                 None)
                if found_ans:
                    final_ans_audio = found_ans
                    if not found_ans.lower().endswith('.mp3'):
                        try:
                            conv_ans_path = os.path.splitext(found_ans)[0] + "_ans_converted.mp3"
                            if not os.path.exists(conv_ans_path):
                                AudioSegment.from_file(found_ans).export(conv_ans_path, format="mp3")
                            final_ans_audio = conv_ans_path
                        except Exception:
                            pass
                    slide.shapes.add_movie(os.path.abspath(final_ans_audio), Inches(-1), Inches(-1), Inches(0.5),
                                           Inches(0.5)).name = "TriviaAudio"
            elif r == 2:
                # 2 Images Side-by-Side + Text at Bottom
                img_q_val = str(row.get('Question'))
                img_path_q = os.path.join("images", f"Round {r}",
                                          img_q_val) if img_q_val and img_q_val.lower() != 'nan' else None

                if img_path_q:
                    place_smart_scaled_image(slide, img_path_q, prs.slide_width * 0.25,
                                             prs.slide_height / 2 - Inches(0.5), max_w=Inches(6.0), max_h=Inches(4.5))
                if path:
                    place_smart_scaled_image(slide, path, prs.slide_width * 0.75, prs.slide_height / 2 - Inches(0.5),
                                             max_w=Inches(6.0), max_h=Inches(4.5))

                add_styled_text(slide, f"Răspuns: {row['Answer']}", 0, Inches(5.8), prs.slide_width, Inches(1.5),
                                font_size=54, bold=True, font_name="Gladiola", force_single_line=False)

            else:
                # Standard Logic for R1, R4, R5 Answer Slides
                add_styled_text(slide, row['Question'], Inches(0.5), Inches(2.2), Inches(6), Inches(2.5), is_qa=True,
                                font_name="Gladiola", force_single_line=False)
                add_styled_text(slide, f"Răspuns: {row['Answer']}", Inches(0.5), Inches(4.7), Inches(6), Inches(2.5),
                                font_size=44, bold=True, is_qa=True, font_name="Gladiola", force_single_line=False)
                if path: place_smart_scaled_image(slide, path, Inches(10), prs.slide_height / 2 + Inches(0.5),
                                                  max_w=Inches(5.0), max_h=Inches(4.5))

        if r in [3, 5]: add_styled_slide("Pauză", "Rezultatele momentului" if r == 3 else "Situația înainte de Pariu")

    # --- WAGER ---
    p_data = df[df['Round'] == 6]
    if not p_data.empty:
        p_row = p_data.iloc[0]
        add_styled_slide("Pariul", "Miza crește...")
        sq = add_styled_slide()
        add_bg(sq)
        add_styled_text(sq, "Pariul: Întrebarea", Inches(0.5), Inches(0.2), Inches(5), Inches(0.8), font_size=44,
                        alignment=PP_ALIGN.LEFT, font_name="Gladiola", color_rgb=(255, 215, 0))
        w_img = glob.glob(os.path.join("images", "Wager", "*.*"))
        if w_img:
            add_styled_text(sq, p_row['Question'], Inches(0.5), Inches(1.0), Inches(6.0), Inches(6.0), bold=True,
                            is_qa=True, font_name="Gladiola", force_single_line=False)
            place_smart_scaled_image(sq, w_img[0], Inches(10), prs.slide_height / 2 + Inches(0.5), max_w=Inches(5.0),
                                     max_h=Inches(4.5))
        else:
            add_styled_text(sq, p_row['Question'], Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0), bold=True,
                            is_qa=True, font_name="Gladiola", force_single_line=False)
        add_styled_slide("Răspunsul", "Momentul Adevărului...")
        sa = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(sa)
        add_styled_text(sa, "Pariul: Răspuns", Inches(0.5), Inches(0.2), Inches(5), Inches(0.8), font_size=44,
                        alignment=PP_ALIGN.LEFT, font_name="Gladiola", color_rgb=(255, 215, 0))
        add_styled_text(sa, f"Răspuns: {p_row['Answer']}", Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0),
                        font_size=54, bold=True, is_qa=True, font_name="Gladiola", force_single_line=False)
        add_styled_slide("Rezultate Finale", "Câștigătorii Transylvania Trivia")

    prs.save("Transylvania_Trivia.pptx")
    return "Transylvania_Trivia.pptx"


# --- 4. APP EXECUTION ---
tab_cmd, tab_score, tab_reveal = st.tabs(["🎮 Trivia Command", "🏆 Scoring Matrix", "📽️ Leaderboard Presentation"])
with tab_cmd:
    st.title("🧛 Trivia Production")
    uploaded = st.file_uploader("Load CSV", type="csv")
    if uploaded and st.button("Generate Slides"):
        path = generate_trivia_slides(pd.read_csv(uploaded))
        with open(path, "rb") as f: st.download_button("Download", f, file_name=path)
with tab_score: manage_scores()
with tab_reveal: show_leaderboard_reveal()